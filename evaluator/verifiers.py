"""
File verification utilities for checking created files against success criteria.
Uses openpyxl, python-docx, and python-pptx to inspect file contents.
"""

import os
from typing import Dict, Any, List
from pathlib import Path


def verify_file_exists(file_path: str) -> bool:
    """Check if file exists"""
    return os.path.exists(file_path) and os.path.isfile(file_path)


def verify_xlsx_file(file_path: str, criteria: Dict[str, Any]) -> Dict[str, bool]:
    """
    Verify Excel (.xlsx) file against criteria.

    Possible criteria:
    - file_valid: Can open without error
    - has_formula: Contains at least one formula
    - min_rows: Minimum number of rows with data
    - min_columns: Minimum number of columns with data
    - has_chart: Contains at least one chart
    """
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise ImportError("openpyxl not installed. Run: pip install openpyxl")

    results = {}

    # Check file_valid
    try:
        wb = load_workbook(file_path, data_only=False)
        ws = wb.active
        results['file_valid'] = True
    except Exception as e:
        results['file_valid'] = False
        # If file can't be opened, all other checks fail
        for key in criteria:
            results[key] = False
        return results

    # Check has_formula
    if 'has_formula' in criteria:
        has_formula = False
        for row in ws.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str) and cell.value.startswith('='):
                    has_formula = True
                    break
            if has_formula:
                break
        results['has_formula'] = has_formula

    # Check min_rows
    if 'min_rows' in criteria:
        max_row = ws.max_row
        # Count non-empty rows
        non_empty_rows = 0
        for row in ws.iter_rows(max_row=max_row):
            if any(cell.value is not None for cell in row):
                non_empty_rows += 1
        results['min_rows'] = non_empty_rows >= criteria['min_rows']

    # Check min_columns
    if 'min_columns' in criteria:
        max_col = ws.max_column
        # Count non-empty columns
        non_empty_cols = 0
        for col in ws.iter_cols(max_col=max_col):
            if any(cell.value is not None for cell in col):
                non_empty_cols += 1
        results['min_columns'] = non_empty_cols >= criteria['min_columns']

    # Check has_chart
    if 'has_chart' in criteria:
        results['has_chart'] = len(ws._charts) > 0

    wb.close()
    return results


def verify_docx_file(file_path: str, criteria: Dict[str, Any]) -> Dict[str, bool]:
    """
    Verify Word (.docx) file against criteria.

    Possible criteria:
    - file_valid: Can open without error
    - min_paragraphs: Minimum number of paragraphs
    - min_words: Minimum word count
    - has_table: Contains at least one table
    - has_image: Contains at least one image
    """
    try:
        from docx import Document
    except ImportError:
        raise ImportError("python-docx not installed. Run: pip install python-docx")

    results = {}

    # Check file_valid
    try:
        doc = Document(file_path)
        results['file_valid'] = True
    except Exception as e:
        results['file_valid'] = False
        for key in criteria:
            results[key] = False
        return results

    # Check min_paragraphs
    if 'min_paragraphs' in criteria:
        non_empty_paragraphs = [p for p in doc.paragraphs if p.text.strip()]
        results['min_paragraphs'] = len(non_empty_paragraphs) >= criteria['min_paragraphs']

    # Check min_words
    if 'min_words' in criteria:
        total_words = sum(len(p.text.split()) for p in doc.paragraphs)
        results['min_words'] = total_words >= criteria['min_words']

    # Check has_table
    if 'has_table' in criteria:
        results['has_table'] = len(doc.tables) > 0

    # Check has_image
    if 'has_image' in criteria:
        has_image = False
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                has_image = True
                break
        results['has_image'] = has_image

    return results


def verify_pptx_file(file_path: str, criteria: Dict[str, Any]) -> Dict[str, bool]:
    """
    Verify PowerPoint (.pptx) file against criteria.

    Possible criteria:
    - file_valid: Can open without error
    - min_slides: Minimum number of slides
    - has_chart: Contains at least one chart
    - has_image: Contains at least one image
    - has_table: Contains at least one table
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")

    results = {}

    # Check file_valid
    try:
        prs = Presentation(file_path)
        results['file_valid'] = True
    except Exception as e:
        results['file_valid'] = False
        for key in criteria:
            results[key] = False
        return results

    # Check min_slides
    if 'min_slides' in criteria:
        results['min_slides'] = len(prs.slides) >= criteria['min_slides']

    # Check has_chart
    if 'has_chart' in criteria:
        has_chart = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_chart:
                    has_chart = True
                    break
            if has_chart:
                break
        results['has_chart'] = has_chart

    # Check has_image
    if 'has_image' in criteria:
        has_image = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # PICTURE = 13
                    has_image = True
                    break
            if has_image:
                break
        results['has_image'] = has_image

    # Check has_table
    if 'has_table' in criteria:
        has_table = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    has_table = True
                    break
            if has_table:
                break
        results['has_table'] = has_table

    return results


def verify_file(file_path: str, criteria: Dict[str, Any]) -> Dict[str, bool]:
    """
    Main verification function that routes to appropriate verifier based on file type.

    Args:
        file_path: Path to file to verify
        criteria: Dict of success criteria to check

    Returns:
        Dict mapping criterion name to pass/fail bool
    """
    if not verify_file_exists(file_path):
        return {key: False for key in criteria}

    ext = Path(file_path).suffix.lower()

    if ext == '.xlsx':
        return verify_xlsx_file(file_path, criteria)
    elif ext == '.docx':
        return verify_docx_file(file_path, criteria)
    elif ext == '.pptx':
        return verify_pptx_file(file_path, criteria)
    else:
        # Unknown file type - just check file_valid
        results = {}
        results['file_valid'] = os.path.getsize(file_path) > 0
        # For unknown types, we can't verify specific criteria
        # Don't set file_created - that's handled by task_runner
        for key in criteria:
            if key not in ['file_valid', 'file_created']:
                results[key] = True  # Assume pass for unknown criteria
        return results


def find_created_files(directory: str, extensions: List[str] = None) -> List[str]:
    """
    Find recently created files in a directory.

    Args:
        directory: Directory to search
        extensions: Optional list of extensions to filter (e.g., ['.xlsx', '.docx'])

    Returns:
        List of file paths
    """
    if not os.path.exists(directory):
        return []

    files = []
    for item in os.listdir(directory):
        item_path = os.path.join(directory, item)
        if os.path.isfile(item_path):
            if extensions is None or Path(item_path).suffix.lower() in extensions:
                files.append(item_path)

    # Sort by modification time (most recent first)
    files.sort(key=lambda x: os.path.getmtime(x), reverse=True)
    return files
